
#########################################################################
### Executable code
#########################################################################

def get_msig(msigterm):
	'''
	get_msig(msigterm)[0] : columns
	get_msig(msigterm)[1] : contents
	'''
	import requests
	import lxml.html as lh
	import pandas as pd
	from bs4 import BeautifulSoup
	url="https://%s.html"%(msigterm)
	response = requests.get(url)
	soup = BeautifulSoup(response.text,'lxml')
	#print("==========================================================\n")
	soup.select('table tr td')[2]
	### extract th as column
	columns = soup.select('table > tr > th')
	columnlist = [column.text for column in columns]
	#print("==========================================================\n")
	### extract td as contents
	contents = soup.select('table > tr')
	dfcontent = []
	alldfcontents = []
	for content in contents:
		tds = content.find_all("td")
		for td in tds :
			dfcontent.append(td.text)
		alldfcontents.append(dfcontent)
		dfcontent = []
	### select columns and contents 
	select_columns = columnlist[:6]
	select_contents = ["".join(dfcons) for dfcons in alldfcontents[1:7]]
	res_select = []
	res_select.append(select_columns)
	res_select.append(select_contents)
	return(res_select)
	#for i in range(len(select_contents)):
		#print(select_columns[i], ":", select_contents[i])

### make ppt slide
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Pt
from pptx.dml.color import RGBColor


presentation = Presentation()

#/home/cytogenbi2/SingleCell/c2
import os
import glob

resf = glob.glob("/home/cytogenbi2/SingleCell/c2/*.csv")
"""
for res in resf:
	with open(res) as f:
		termlist = [fline.split(",")[0].replace('"','') for fline in f.readlines()]
"""
clusterfile = resf[2]
with open(clusterfile) as f:
	termlist = [fline.split(",")[0].replace('"','') for fline in f.readlines()][:30]

linenum=0
for msigterm in termlist:
	linenum+=1
	print("Term number = %s"%linenum)
	if len(msigterm) < 1 :
		continue
	select_columns = get_msig(msigterm)[0]
	select_contents = get_msig(msigterm)[1]
	### title only slide
	title_only_slide_layout = presentation.slide_layouts[5]
	slide = presentation.slides.add_slide(title_only_slide_layout)
	### set slide title
	shapes = slide.shapes
	shapes.title.text = msigterm.replace("_"," ")
	### set table rows and cols
	table = shapes.add_table(rows=6, cols=2, left=Inches(0.0), top=Inches(2.0), width=Inches(10.0), height=Inches(3.0)).table
	table.columns[0].width = Inches(4.0)
	table.columns[1].width = Inches(6.0)
	for i in range(len(select_contents)):
		table.cell(i,0).text = select_columns[i]
		table.cell(i,1).text = select_contents[i]
		print(i)
		#print("table.cell(0,%s).text = select_columns[%s]"%(i,i))
		#print("table.cell(1,%s).text = select_contents[%s]"%(i,i))
	for i in range(len(select_contents)):
		for j in (0,1):
			cell = table.cell(i,j)
			para = cell.text_frame.paragraphs[0]
			para.font.size = Pt(10)
			para.font.name = 'Comic Sans MS'
			para.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT


presentation.save('/home/cytogenbi2/SingleCell/c2/%s.pptx'%(clusterfile.split("/")[-1].split("_result")[0]))




#########################################################################
### Executable code
#########################################################################

#/home/cytogenbi2/SingleCell/c2


import requests
import lxml.html as lh
import pandas as pd
from bs4 import BeautifulSoup

url="https://.html"

msigterm = ''
response = requests.get(url)

soup = BeautifulSoup(response.text,'lxml')

#print("==========================================================\n")
soup.select('table tr td')[2]

### extract th as column
columns = soup.select('table > tr > th')
columnlist = [column.text for column in columns]


#print("==========================================================\n")

### extract td as contents
contents = soup.select('table > tr')
dfcontent = []
alldfcontents = []

for content in contents:
	tds = content.find_all("td")
	for td in tds :
		dfcontent.append(td.text)
	alldfcontents.append(dfcontent)
	dfcontent = []

### select columns and contents 
select_columns = columnlist[:6]
select_contents = ["".join(dfcons) for dfcons in alldfcontents[1:7]]

for i in range(len(select_contents)):
	print(select_columns[i], ":", select_contents[i])

### make ppt slide
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Pt
from pptx.dml.color import RGBColor


presentation = Presentation()

### title only slide
title_only_slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(title_only_slide_layout)

### set slide title
shapes = slide.shapes
shapes.title.text = msigterm.replace("_"," ")

### set table rows and cols
table = shapes.add_table(rows=6, cols=2, left=Inches(0.0), top=Inches(2.0), width=Inches(10.0), height=Inches(3.0)).table

table.columns[0].width = Inches(4.0)
table.columns[1].width = Inches(6.0)

for i in range(len(select_contents)):
	table.cell(i,0).text = select_columns[i]
	table.cell(i,1).text = select_contents[i]
	print(i)
	#print("table.cell(0,%s).text = select_columns[%s]"%(i,i))
	#print("table.cell(1,%s).text = select_contents[%s]"%(i,i))


for i in range(len(select_contents)):
	for j in (0,1):
		cell = table.cell(i,j)
		para = cell.text_frame.paragraphs[0]
		para.font.size = Pt(10)
		para.font.name = 'Comic Sans MS'
		para.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT


presentation.save('/home/cytogenbi2/SingleCell/c2/test.pptx')
#########################################################################
### Code practice
#########################################################################


#df = pd.DataFrame(columns=select_columns, data=select_contents)


stable = soup.find('table',{"class":"lists4"})

all_tr_stable = stable.find_all("tr")
print(all_tr_stable)

### remove tag + str format change ###
tdict={}
for tr in all_tr_stable:
	tr_th = tr.find_all("th")
	thlist=[]
	for th in tr_th:
		thlist.append(str(th.text).replace("\n",""))		
	tr_td = tr.find_all("td")
	tdlist=[]
	for td in tr_td:
		tdlist.append(str(td.text).replace("\n",""))
	tdict["\t".join(thlist)] = "\t".join(tdlist)

### str format change ###
tdict={}
for tr in all_tr_stable:
	tr_th = tr.find_all("th")
	thlist=[]
	for th in tr_th:
		thlist.append(str(th).replace("\n",""))		
	tr_td = tr.find_all("td")
	tdlist=[]
	for td in tr_td:
		tdlist.append(str(td).replace("\n",""))
	tdict["\t".join(thlist)] = "\t".join(tdlist)


for tr in all_tr_stable:
	tr_th = tr.find_all("th")
	thlist=[]
	for th in tr_th:
		valign = th.select('td > valign')
		print(valign)
	break
	tr_td = tr.find_all("td")
	tdlist=[]
	for td in tr_td:
		tdlist.append(str(td).replace("\n",""))
	tdict["\t".join(thlist)] = "\t".join(tdlist)
	
	
AttributeError: 'str' object has no attribute 'select'



for tr in all_tr_stable:
	#tr_th = tr.find_all("th")
	#print(tr_th)	
	tr_td = tr.find_all("td")
	print(tr_td)
	
	

	
	
all_th_stable = stable.find_all("th")
print(all_th_stable)

all_td_stable = stable.find_all("td")
print(all_td_stable)



#########################################################################

for points in soup.find_all('div',{"class":"user-level"}):
    point = str(points.text)
    print(point)




